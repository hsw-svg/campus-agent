"""Render a normalised slide_deck dict to a styled .pptx binary via python-pptx.

Uses built-in themed templates from pptx_templates/ directory.
Falls back to blank Presentation if no template is available.
"""

from __future__ import annotations

import io
from dataclasses import dataclass
from pathlib import Path
from typing import Any

_TEMPLATE_DIR = Path(__file__).parent / "pptx_templates"

# Theme selection keywords — first match wins
_THEME_KEYWORDS: tuple[tuple[str, str], ...] = (
    ("编程", "tech_purple"),
    ("代码", "tech_purple"),
    ("算法", "tech_purple"),
    ("Python", "tech_purple"),
    ("python", "tech_purple"),
    ("Java", "tech_purple"),
    ("javascript", "tech_purple"),
    ("数据结构", "tech_purple"),
    ("机器学习", "tech_purple"),
    ("深度学习", "tech_purple"),
    ("生物", "vibrant_green"),
    ("化学", "vibrant_green"),
    ("物理", "vibrant_green"),
    ("实验", "vibrant_green"),
    ("自然", "vibrant_green"),
    ("生态", "vibrant_green"),
    ("医学", "vibrant_green"),
    ("文学", "warm_orange"),
    ("历史", "warm_orange"),
    ("哲学", "warm_orange"),
    ("艺术", "warm_orange"),
    ("文化", "warm_orange"),
    ("社会", "warm_orange"),
    ("教育", "warm_orange"),
    ("心理", "warm_orange"),
)

# Theme metadata for font styling
_THEME_META: dict[str, dict[str, str]] = {
    "academic_blue": {
        "cover_title": "FFFFFF",
        "cover_sub": "B0C4DE",
        "title_clr": "1A3A5C",
        "text_clr": "2C3E50",
    },
    "vibrant_green": {
        "cover_title": "FFFFFF",
        "cover_sub": "C8E6C9",
        "title_clr": "1B5E20",
        "text_clr": "33691E",
    },
    "warm_orange": {
        "cover_title": "FFFFFF",
        "cover_sub": "FFCCBC",
        "title_clr": "BF360C",
        "text_clr": "4E342E",
    },
    "tech_purple": {
        "cover_title": "FFFFFF",
        "cover_sub": "D1C4E9",
        "title_clr": "4A148C",
        "text_clr": "311B92",
    },
}

_DEFAULT_THEME = "academic_blue"


@dataclass(frozen=True)
class ExportedBinaryArtifact:
    content: bytes
    media_type: str
    extension: str


def _select_theme(data: dict[str, Any]) -> str:
    """Pick a template theme based on the deck topic + context signals."""
    text_parts: list[str] = []
    text_parts.append(str(data.get("topic") or ""))
    signals = data.get("context_signals") or {}
    for key in ("learning_analysis", "classroom_summary", "grading", "weak_points"):
        val = signals.get(key)
        if isinstance(val, list):
            text_parts.extend(str(v) for v in val)
        elif val:
            text_parts.append(str(val))
    text = " ".join(text_parts)

    for keyword, theme_key in _THEME_KEYWORDS:
        if keyword in text:
            return theme_key
    return _DEFAULT_THEME


def _load_template(theme_key: str):
    """Load a themed template, or return None if unavailable."""
    from pptx import Presentation

    path = _TEMPLATE_DIR / f"{theme_key}.pptx"
    if path.exists():
        return Presentation(str(path))
    # Fallback to default theme
    fallback = _TEMPLATE_DIR / f"{_DEFAULT_THEME}.pptx"
    if fallback.exists():
        return Presentation(str(fallback))
    # Last resort: blank presentation
    return Presentation()


def _hex_rgb(hex_color: str):
    from pptx.dml.color import RGBColor

    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _apply_font_color(text_frame, color_hex: str) -> None:
    """Set font color on all runs in a text frame."""
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = _hex_rgb(color_hex)
            run.font.name = "Microsoft YaHei"


class SlideDeckPptxSkill:
    id = "slide_deck_pptx"
    version = "2"
    input_type = "dict"
    output_type = "ExportedBinaryArtifact"
    error_codes = ()
    has_side_effects = False
    can_access_workspace = False

    def run(self, data: dict[str, Any]) -> ExportedBinaryArtifact:
        theme_key = _select_theme(data)
        meta = _THEME_META.get(theme_key, _THEME_META[_DEFAULT_THEME])
        presentation = _load_template(theme_key)
        layouts = presentation.slide_layouts

        title_layout = layouts[0]
        content_layout = layouts[1] if len(layouts) > 1 else layouts[0]
        two_col_layout = layouts[5] if len(layouts) > 5 else (
            layouts[3] if len(layouts) > 3 else content_layout
        )

        topic = str(data.get("topic") or "未命名课件").strip()
        audience = str(data.get("audience") or "").strip()
        objective = str(data.get("objective") or "").strip()

        # Cover slide
        cover = presentation.slides.add_slide(title_layout)
        _set_title(cover, topic)
        _apply_font_color(cover.shapes.title.text_frame, meta["cover_title"])
        subtitle_bits: list[str] = []
        if audience:
            subtitle_bits.append(audience)
        if objective:
            subtitle_bits.append(objective)
        if subtitle_bits:
            _set_subtitle(cover, " · ".join(subtitle_bits))
            _apply_font_color(
                _find_placeholder_text_frame(cover, 1), meta["cover_sub"]
            )

        # Content slides
        for slide in data.get("slides") or []:
            layout_key = slide.get("layout") or "bullets"
            if layout_key == "title":
                pptx_slide = presentation.slides.add_slide(title_layout)
                _set_title(pptx_slide, str(slide.get("title") or ""))
                _apply_font_color(pptx_slide.shapes.title.text_frame, meta["cover_title"])
                if slide.get("subtitle"):
                    _set_subtitle(pptx_slide, str(slide["subtitle"]))
                    _apply_font_color(
                        _find_placeholder_text_frame(pptx_slide, 1), meta["cover_sub"]
                    )
            elif layout_key == "two_column" and slide.get("columns"):
                pptx_slide = presentation.slides.add_slide(two_col_layout)
                _set_title(pptx_slide, str(slide.get("title") or ""))
                _apply_font_color(pptx_slide.shapes.title.text_frame, meta["title_clr"])
                _fill_two_column(pptx_slide, slide.get("columns") or [])
                _apply_font_color(
                    _find_body_placeholder(pptx_slide).text_frame, meta["text_clr"]
                ) if _find_body_placeholder(pptx_slide) else None
            else:
                pptx_slide = presentation.slides.add_slide(content_layout)
                _set_title(pptx_slide, str(slide.get("title") or ""))
                _apply_font_color(pptx_slide.shapes.title.text_frame, meta["title_clr"])
                bullets = list(slide.get("bullets") or [])
                if layout_key == "callout" and bullets:
                    bullets = [f"重点：{bullets[0]}"] + bullets[1:]
                if slide.get("key_points"):
                    bullets = bullets + [
                        "重点：" + "、".join(slide["key_points"])
                    ]
                _fill_body(pptx_slide, bullets)
                body = _find_body_placeholder(pptx_slide)
                if body is not None:
                    _apply_font_color(body.text_frame, meta["text_clr"])

            notes = str(slide.get("notes") or "").strip()
            media = list(slide.get("media") or [])
            media_notes = _build_media_notes(media)
            combined_notes = "\n".join(filter(None, [notes, media_notes])).strip()
            if combined_notes:
                pptx_slide.notes_slide.notes_text_frame.text = combined_notes

        buffer = io.BytesIO()
        presentation.save(buffer)
        return ExportedBinaryArtifact(
            content=buffer.getvalue(),
            media_type=(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
            extension="pptx",
        )


def _set_title(slide, text: str) -> None:
    if slide.shapes.title is not None:
        slide.shapes.title.text = text


def _set_subtitle(slide, text: str) -> None:
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            shape.text = text
            return


def _find_placeholder_text_frame(slide, idx: int):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape.text_frame
    return None


def _fill_body(slide, bullets: list[str]) -> None:
    body = _find_body_placeholder(slide)
    if body is None or not bullets:
        return
    text_frame = body.text_frame
    text_frame.text = bullets[0]
    for extra in bullets[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = extra


def _fill_two_column(slide, columns: list[dict[str, Any]]) -> None:
    placeholders = [
        shape
        for shape in slide.placeholders
        if shape.placeholder_format.idx not in (0,)
    ]
    for idx, column in enumerate(columns[:2]):
        if idx >= len(placeholders):
            break
        placeholder = placeholders[idx]
        header = column.get("title") or ""
        bullets = list(column.get("bullets") or [])
        text_frame = placeholder.text_frame
        first_line = header or (bullets.pop(0) if bullets else "")
        text_frame.text = first_line
        for bullet in bullets:
            paragraph = text_frame.add_paragraph()
            paragraph.text = bullet


def _find_body_placeholder(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            return shape
    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 0:
            return shape
    return None


_MEDIA_KIND_LABELS: dict[str, str] = {
    "image": "图片",
    "video": "视频",
    "gif": "动图",
    "embed": "嵌入",
    "animation": "动画",
}

_PLACEMENT_LABELS: dict[str, str] = {
    "top": "顶部",
    "right": "右侧",
    "left": "左侧",
    "background": "背景",
    "inline": "内嵌",
}


def _build_media_notes(media: list[dict[str, Any]]) -> str:
    if not media:
        return ""
    lines: list[str] = ["[多媒体建议]"]
    for item in media:
        kind = str(item.get("kind") or "素材")
        kind_label = _MEDIA_KIND_LABELS.get(kind, kind)
        title = str(item.get("title") or item.get("alt") or "未命名素材")
        url = str(item.get("url") or "").strip()
        placement = str(item.get("placement") or "inline")
        placement_label = _PLACEMENT_LABELS.get(placement, placement)
        line = f"• {kind_label}: {title}（{placement_label}）"
        if url:
            line += f"\n  链接: {url}"
        lines.append(line)
    return "\n".join(lines)
