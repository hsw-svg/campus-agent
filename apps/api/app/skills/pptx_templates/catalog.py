"""Controlled catalog for course slide-deck templates and editable frames."""

from __future__ import annotations

import json
from dataclasses import dataclass
from functools import lru_cache
from pathlib import Path
from typing import Any

from app.core.errors import AppError

TEMPLATE_DIR = Path(__file__).parent
MANIFEST_DIR = TEMPLATE_DIR / "manifests"
DEFAULT_TEMPLATE_ID = "ai_tech"
_MANIFEST_NAMES = ("ai_tech.json", "business_plan.json")


@dataclass(frozen=True)
class TextTargetSpec:
    shape_id: int
    action: str = "rewrite"
    field: str | None = None
    max_chars: int | None = None


@dataclass(frozen=True)
class PptxFrameSpec:
    id: str
    source_slide: int
    layouts: tuple[str, ...]
    max_items: int
    targets: tuple[TextTargetSpec, ...]
    keep_shape_ids: tuple[int, ...]


@dataclass(frozen=True)
class PptxTemplateSpec:
    id: str
    display_name: str
    description: str
    file_name: str
    license_scope: str
    aliases: tuple[str, ...]
    frames: tuple[PptxFrameSpec, ...]

    @property
    def path(self) -> Path:
        return TEMPLATE_DIR / self.file_name

    def frames_for_layout(self, layout: str) -> tuple[PptxFrameSpec, ...]:
        return tuple(frame for frame in self.frames if layout in frame.layouts)

    def frame(self, frame_id: str) -> PptxFrameSpec:
        for frame in self.frames:
            if frame.id == frame_id:
                return frame
        raise KeyError(frame_id)


def _manifest_error(message: str, **details: Any) -> AppError:
    return AppError(
        code="pptx_template_manifest_invalid",
        message=message,
        status_code=500,
        details=details or None,
    )


def _load_manifest(path: Path) -> PptxTemplateSpec:
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
        frames = tuple(
            PptxFrameSpec(
                id=str(frame["id"]),
                source_slide=int(frame["source_slide"]),
                layouts=tuple(str(value) for value in frame["layouts"]),
                max_items=int(frame.get("max_items") or 0),
                targets=tuple(
                    TextTargetSpec(
                        shape_id=int(target["shape_id"]),
                        action=str(target.get("action") or "rewrite"),
                        field=(str(target["field"]) if target.get("field") else None),
                        max_chars=(
                            int(target["max_chars"])
                            if target.get("max_chars") is not None
                            else None
                        ),
                    )
                    for target in frame.get("targets", [])
                ),
                keep_shape_ids=tuple(int(value) for value in frame.get("keep_shape_ids", [])),
            )
            for frame in payload["frames"]
        )
        spec = PptxTemplateSpec(
            id=str(payload["id"]),
            display_name=str(payload["display_name"]),
            description=str(payload["description"]),
            file_name=str(payload["file_name"]),
            license_scope=str(payload["license_scope"]),
            aliases=tuple(str(value) for value in payload.get("aliases", [])),
            frames=frames,
        )
    except (KeyError, TypeError, ValueError, json.JSONDecodeError) as error:
        raise _manifest_error(
            "PPTX template manifest is invalid.",
            manifest=str(path),
            error=str(error),
        ) from error
    if Path(spec.file_name).name != spec.file_name or not spec.file_name.endswith(".pptx"):
        raise _manifest_error(
            "PPTX template file name must be a local .pptx name.",
            template_id=spec.id,
            file_name=spec.file_name,
        )
    return spec


@lru_cache
def template_catalog() -> tuple[PptxTemplateSpec, ...]:
    specs = tuple(_load_manifest(MANIFEST_DIR / name) for name in _MANIFEST_NAMES)
    ids = [spec.id for spec in specs]
    if len(ids) != len(set(ids)) or DEFAULT_TEMPLATE_ID not in ids:
        raise _manifest_error("PPTX template ids must be unique and include the default.")
    return specs


def get_template(template_id: str | None) -> PptxTemplateSpec:
    requested = (template_id or DEFAULT_TEMPLATE_ID).strip()
    for spec in template_catalog():
        if spec.id == requested:
            return spec
    return get_template(DEFAULT_TEMPLATE_ID)


def is_known_template_id(template_id: Any) -> bool:
    value = str(template_id or "").strip()
    return any(spec.id == value for spec in template_catalog())


def find_explicit_template_id(text: str) -> str | None:
    lowered = (text or "").casefold()
    matches: list[tuple[int, str]] = []
    for spec in template_catalog():
        for alias in (spec.display_name, *spec.aliases):
            if alias.casefold() in lowered:
                matches.append((len(alias), spec.id))
    if not matches:
        return None
    matches.sort(reverse=True)
    return matches[0][1]


def prompt_template_catalog() -> list[dict[str, str]]:
    return [
        {"id": spec.id, "name": spec.display_name, "description": spec.description}
        for spec in template_catalog()
    ]


def template_metadata(template_id: str, selection_source: str) -> dict[str, str]:
    spec = get_template(template_id)
    return {
        "template_id": spec.id,
        "template_name": spec.display_name,
        "template_selection_source": selection_source,
        "template_license_scope": spec.license_scope,
    }


def validate_template_manifest(spec: PptxTemplateSpec) -> None:
    """Fail fast when a template edit invalidates a registered object map."""
    if not spec.path.exists():
        raise _manifest_error(
            "PPTX template file is missing.",
            template_id=spec.id,
            path=str(spec.path),
        )

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(str(spec.path))

    def text_shape_ids(shapes) -> set[int]:
        ids: set[int] = set()
        for shape in shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                ids.add(shape.shape_id)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                ids.update(text_shape_ids(shape.shapes))
        return ids

    for frame in spec.frames:
        if frame.source_slide < 1 or frame.source_slide > len(presentation.slides):
            raise _manifest_error(
                "PPTX frame references a missing source slide.",
                template_id=spec.id,
                frame_id=frame.id,
                source_slide=frame.source_slide,
            )
        if not frame.layouts:
            raise _manifest_error(
                "PPTX frame must support at least one logical layout.",
                template_id=spec.id,
                frame_id=frame.id,
            )
        target_ids = {target.shape_id for target in frame.targets}
        classified_ids = target_ids | set(frame.keep_shape_ids)
        actual_ids = text_shape_ids(presentation.slides[frame.source_slide - 1].shapes)
        if classified_ids != actual_ids:
            raise _manifest_error(
                "PPTX frame text objects are not fully classified.",
                template_id=spec.id,
                frame_id=frame.id,
                missing=sorted(actual_ids - classified_ids),
                unknown=sorted(classified_ids - actual_ids),
            )
        if len(target_ids) != len(frame.targets):
            raise _manifest_error(
                "PPTX frame contains duplicate target shape ids.",
                template_id=spec.id,
                frame_id=frame.id,
            )
        for target in frame.targets:
            if target.action not in {"rewrite", "delete"}:
                raise _manifest_error(
                    "PPTX target action is unsupported.",
                    template_id=spec.id,
                    frame_id=frame.id,
                    action=target.action,
                )
            if target.action == "rewrite" and not target.field:
                raise _manifest_error(
                    "PPTX rewrite target must declare a field.",
                    template_id=spec.id,
                    frame_id=frame.id,
                    shape_id=target.shape_id,
                )
