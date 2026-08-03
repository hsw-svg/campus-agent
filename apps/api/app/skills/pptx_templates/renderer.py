"""Render structured slide-deck data by cloning registered template slides."""

from __future__ import annotations

import copy
import io
import re
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from lxml import etree

from app.core.errors import AppError
from app.skills.pptx_templates.catalog import (
    PptxFrameSpec,
    PptxTemplateSpec,
    TextTargetSpec,
    validate_template_manifest,
)

_PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
_CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
_SLIDE_REL_TYPE = f"{_REL_NS}/slide"
_SLIDE_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
)
_SLIDE_PATH_RE = re.compile(r"^ppt/slides/slide\d+\.xml$")
_SLIDE_RELS_PATH_RE = re.compile(r"^ppt/slides/_rels/slide\d+\.xml\.rels$")


@dataclass(frozen=True)
class PlannedSlide:
    frame: PptxFrameSpec
    data: dict[str, Any]


def _render_error(message: str, **details: Any) -> AppError:
    return AppError(
        code="pptx_template_render_failed",
        message=message,
        status_code=500,
        details=details or None,
    )


def _synthesised_cover(data: dict[str, Any]) -> dict[str, Any]:
    audience = str(data.get("audience") or "").strip()
    objective = str(data.get("objective") or "").strip()
    return {
        "index": 0,
        "layout": "title",
        "title": str(data.get("topic") or "未命名课件").strip(),
        "subtitle": " · ".join(value for value in (audience, objective) if value),
        "bullets": [],
        "columns": [],
        "key_points": [],
        "notes": "",
        "media": [],
    }


def _item_count(slide: dict[str, Any]) -> int:
    if str(slide.get("layout") or "") == "two_column":
        return sum(len(column.get("bullets") or []) for column in slide.get("columns") or [])
    return len(slide.get("bullets") or [])


def plan_template_slides(
    spec: PptxTemplateSpec,
    data: dict[str, Any],
) -> tuple[PlannedSlide, ...]:
    slides = [dict(slide) for slide in data.get("slides") or []]
    if not slides or str(slides[0].get("layout") or "") != "title":
        slides.insert(0, _synthesised_cover(data))

    planned: list[PlannedSlide] = []
    previous_frame_id: str | None = None
    for slide in slides:
        layout = str(slide.get("layout") or "bullets").strip().lower()
        candidates = list(spec.frames_for_layout(layout))
        if not candidates:
            candidates = list(spec.frames_for_layout("bullets"))
        if not candidates:
            raise _render_error(
                "PPTX template has no frame for the requested layout.",
                template_id=spec.id,
                layout=layout,
            )
        count = _item_count(slide)
        fitting = [frame for frame in candidates if frame.max_items >= count]
        pool = fitting or candidates
        if previous_frame_id and len(pool) > 1:
            non_repeating = [frame for frame in pool if frame.id != previous_frame_id]
            if non_repeating:
                pool = non_repeating
        frame = min(pool, key=lambda item: (max(item.max_items - count, 0), item.id))
        planned.append(PlannedSlide(frame=frame, data=slide))
        previous_frame_id = frame.id
    return tuple(planned)


def _read_template_entries(template_path: Path) -> dict[str, bytes]:
    with zipfile.ZipFile(template_path) as archive:
        return {name: archive.read(name) for name in archive.namelist() if not name.endswith("/")}


def _clean_slide_relationships(xml: bytes) -> bytes:
    root = etree.fromstring(xml)
    for rel in list(root):
        rel_type = str(rel.get("Type") or "")
        if rel_type.endswith("/notesSlide") or rel_type.endswith("/comments"):
            root.remove(rel)
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def _clone_template_slides(
    template_path: Path,
    source_slide_numbers: tuple[int, ...],
) -> io.BytesIO:
    entries = _read_template_entries(template_path)
    try:
        source_slides = {
            number: entries[f"ppt/slides/slide{number}.xml"]
            for number in set(source_slide_numbers)
        }
    except KeyError as error:
        raise _render_error(
            "PPTX template source slide is missing.",
            template=str(template_path),
            missing=str(error),
        ) from error
    source_rels: dict[int, bytes | None] = {}
    for number in set(source_slide_numbers):
        source_rels[number] = entries.get(
            f"ppt/slides/_rels/slide{number}.xml.rels"
        )

    presentation_root = etree.fromstring(entries["ppt/presentation.xml"])
    slide_id_list = presentation_root.find(f"{{{_PML_NS}}}sldIdLst")
    if slide_id_list is None:
        raise _render_error("PPTX template presentation has no slide id list.")
    for child in list(slide_id_list):
        slide_id_list.remove(child)

    presentation_rels = etree.fromstring(entries["ppt/_rels/presentation.xml.rels"])
    for rel in list(presentation_rels):
        if rel.get("Type") == _SLIDE_REL_TYPE:
            presentation_rels.remove(rel)

    for index in range(1, len(source_slide_numbers) + 1):
        relationship_id = f"rIdGeneratedSlide{index}"
        rel = etree.SubElement(presentation_rels, f"{{{_PKG_REL_NS}}}Relationship")
        rel.set("Id", relationship_id)
        rel.set("Type", _SLIDE_REL_TYPE)
        rel.set("Target", f"slides/slide{index}.xml")
        slide_id = etree.SubElement(slide_id_list, f"{{{_PML_NS}}}sldId")
        slide_id.set("id", str(255 + index))
        slide_id.set(f"{{{_REL_NS}}}id", relationship_id)

    content_types = etree.fromstring(entries["[Content_Types].xml"])
    for override in list(content_types):
        part_name = str(override.get("PartName") or "")
        if re.fullmatch(r"/ppt/slides/slide\d+\.xml", part_name):
            content_types.remove(override)
    for index in range(1, len(source_slide_numbers) + 1):
        override = etree.SubElement(content_types, f"{{{_CONTENT_TYPES_NS}}}Override")
        override.set("PartName", f"/ppt/slides/slide{index}.xml")
        override.set("ContentType", _SLIDE_CONTENT_TYPE)

    output_entries = {
        name: content
        for name, content in entries.items()
        if not _SLIDE_PATH_RE.fullmatch(name) and not _SLIDE_RELS_PATH_RE.fullmatch(name)
    }
    output_entries["ppt/presentation.xml"] = etree.tostring(
        presentation_root, xml_declaration=True, encoding="UTF-8", standalone=True
    )
    output_entries["ppt/_rels/presentation.xml.rels"] = etree.tostring(
        presentation_rels, xml_declaration=True, encoding="UTF-8", standalone=True
    )
    output_entries["[Content_Types].xml"] = etree.tostring(
        content_types, xml_declaration=True, encoding="UTF-8", standalone=True
    )

    for index, source_number in enumerate(source_slide_numbers, start=1):
        output_entries[f"ppt/slides/slide{index}.xml"] = source_slides[source_number]
        rels = source_rels[source_number]
        if rels:
            output_entries[f"ppt/slides/_rels/slide{index}.xml.rels"] = (
                _clean_slide_relationships(rels)
            )

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, content in output_entries.items():
            archive.writestr(name, content)
    buffer.seek(0)
    return buffer


def _iter_shapes(shapes):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _shape_by_id(slide, shape_id: int):
    for shape in _iter_shapes(slide.shapes):
        if shape.shape_id == shape_id:
            return shape
    return None


def _delete_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def _replace_paragraph_text(paragraph, text: str) -> None:
    runs = paragraph.runs
    if runs:
        runs[0].text = text
        for run in runs[1:]:
            run.text = ""
    else:
        paragraph.add_run().text = text


def _set_text_preserving_style(shape, text: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        raise _render_error(
            "PPTX template target does not have a text frame.",
            shape_id=shape.shape_id,
            shape_name=shape.name,
        )
    frame = shape.text_frame
    lines = text.splitlines() or [""]
    paragraphs = list(frame.paragraphs)
    template_paragraph = paragraphs[0]
    while len(paragraphs) < len(lines):
        paragraph = frame.add_paragraph()
        if template_paragraph._p.pPr is not None:
            paragraph._p.insert(0, copy.deepcopy(template_paragraph._p.pPr))
        paragraphs = list(frame.paragraphs)
    for index, paragraph in enumerate(paragraphs):
        _replace_paragraph_text(paragraph, lines[index] if index < len(lines) else "")


def _clip(value: str, max_chars: int | None) -> tuple[str, bool]:
    cleaned = "\n".join(
        " ".join(line.split())
        for line in str(value or "").splitlines()
        if line.strip()
    )
    if max_chars is None or len(cleaned) <= max_chars:
        return cleaned, False
    if max_chars <= 1:
        return cleaned[:max_chars], True
    return cleaned[: max_chars - 1].rstrip() + "…", True


def _split_item(value: str) -> tuple[str, str]:
    cleaned = " ".join(str(value or "").split())
    for separator in ("：", ":", "—", "-"):
        if separator in cleaned:
            title, body = cleaned.split(separator, 1)
            if title.strip() and body.strip():
                return title.strip(), body.strip()
    if len(cleaned) <= 12:
        return cleaned, ""
    return cleaned[:12].rstrip("，。； "), cleaned[12:].lstrip("，。； ")


def _slide_fields(
    deck: dict[str, Any],
    slide: dict[str, Any],
) -> dict[str, str]:
    bullets = [str(value).strip() for value in slide.get("bullets") or [] if str(value).strip()]
    key_points = [
        str(value).strip() for value in slide.get("key_points") or [] if str(value).strip()
    ]
    if not bullets and key_points:
        bullets = key_points
    topic = str(deck.get("topic") or "未命名课件").strip()
    title = str(slide.get("title") or topic).strip()
    subtitle = str(slide.get("subtitle") or deck.get("objective") or "").strip()
    fields: dict[str, str] = {
        "title": title,
        "subtitle": subtitle,
        "eyebrow": "COURSE PRESENTATION",
        "english_label": "Course Presentation",
        "body": "\n".join(bullets),
    }

    for index in range(1, 7):
        value = bullets[index - 1] if index <= len(bullets) else ""
        item_title, item_body = _split_item(value)
        fields[f"item_{index}"] = "\n".join(
            part for part in (item_title, item_body) if part
        )
        fields[f"item_{index}_title"] = item_title
        fields[f"item_{index}_body"] = item_body

    columns = list(slide.get("columns") or [])[:2]
    for index in range(1, 3):
        column = columns[index - 1] if index <= len(columns) else {}
        column_bullets = [
            str(value).strip()
            for value in column.get("bullets") or []
            if str(value).strip()
        ]
        fields[f"column_{index}_title"] = str(column.get("title") or "").strip()
        fields[f"column_{index}_body"] = "\n".join(column_bullets)
    if columns:
        flattened = [
            value
            for column in columns
            for value in column.get("bullets") or []
            if str(value).strip()
        ]
        for index, value in enumerate(flattened[:4], start=1):
            item_title, item_body = _split_item(str(value))
            fields[f"item_{index}"] = "\n".join(
                part for part in (item_title, item_body) if part
            )
            fields[f"item_{index}_title"] = item_title
            fields[f"item_{index}_body"] = item_body
    return fields


def _target_value(
    target: TextTargetSpec,
    fields: dict[str, str],
) -> tuple[str, str | None]:
    full = fields.get(target.field or "", "")
    clipped, changed = _clip(full, target.max_chars)
    return clipped, full if changed else None


def _media_notes(media: list[dict[str, Any]]) -> str:
    if not media:
        return ""
    lines = ["[多媒体建议]"]
    for item in media:
        title = str(item.get("title") or item.get("alt") or "未命名素材")
        url = str(item.get("url") or "").strip()
        lines.append(f"• {title}" + (f"\n  链接: {url}" if url else ""))
    return "\n".join(lines)


def _fill_slide(
    slide,
    planned: PlannedSlide,
    deck: dict[str, Any],
) -> None:
    fields = _slide_fields(deck, planned.data)
    overflow_notes: list[str] = []
    for target in planned.frame.targets:
        shape = _shape_by_id(slide, target.shape_id)
        if shape is None:
            raise _render_error(
                "PPTX template target shape is missing.",
                frame_id=planned.frame.id,
                shape_id=target.shape_id,
            )
        if target.action == "delete":
            _delete_shape(shape)
            continue
        value, overflow = _target_value(target, fields)
        if not value and getattr(shape, "is_placeholder", False):
            _delete_shape(shape)
        else:
            _set_text_preserving_style(shape, value)
        if overflow:
            overflow_notes.append(f"{target.field}: {overflow}")

    notes = str(planned.data.get("notes") or "").strip()
    media = _media_notes(list(planned.data.get("media") or []))
    overflow = "\n".join(["[版面压缩的完整内容]", *overflow_notes]) if overflow_notes else ""
    combined = "\n".join(part for part in (notes, media, overflow) if part).strip()
    if combined:
        slide.notes_slide.notes_text_frame.text = combined


def render_template_deck(
    spec: PptxTemplateSpec,
    data: dict[str, Any],
) -> bytes:
    validate_template_manifest(spec)
    planned = plan_template_slides(spec, data)
    starter = _clone_template_slides(
        spec.path,
        tuple(item.frame.source_slide for item in planned),
    )

    from pptx import Presentation

    try:
        presentation = Presentation(starter)
    except Exception as error:  # noqa: BLE001
        raise _render_error(
            "Cloned PPTX template could not be opened.",
            template_id=spec.id,
            error=str(error),
        ) from error
    if len(presentation.slides) != len(planned):
        raise _render_error(
            "Cloned PPTX template has an unexpected slide count.",
            template_id=spec.id,
            expected=len(planned),
            actual=len(presentation.slides),
        )
    for slide, item in zip(presentation.slides, planned, strict=True):
        _fill_slide(slide, item, data)

    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()
