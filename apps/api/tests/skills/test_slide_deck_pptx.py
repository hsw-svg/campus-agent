"""Regression coverage for the slide_deck -> .pptx render skill."""

import io
import zipfile

import pytest
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.skills.slide_deck_json import SlideDeckJsonSkill
from app.skills.slide_deck_pptx import SlideDeckPptxSkill
from app.skills.pptx_templates.catalog import (
    find_explicit_template_id,
    get_template,
    template_catalog,
    validate_template_manifest,
)
from app.skills.pptx_templates.renderer import plan_template_slides


SAMPLE_DECK = {
    "topic": "Python 切片与元组",
    "audience": "计算机学院大二",
    "objective": "掌握切片语法与元组不可变性",
    "duration_minutes": 45,
    "context_signals": {
        "learning_analysis": "上一次作业平均 82，切片题正确率 61%。",
        "weak_points": ["切片负步长"],
        "classroom_summary": "",
        "grading": "",
        "job_skill_focus": ["中级 Python 岗位常考"],
        "industry_updates": [
            {"title": "案例A", "url": "https://example.com/a", "snippet": "..."},
        ],
    },
    "slides": [
        {
            "index": 1,
            "layout": "title",
            "title": "Python 切片与元组",
            "subtitle": "面向大二·45 分钟",
            "bullets": [],
            "notes": "开场 3 秒抢答",
            "key_points": [],
            "citations": [],
            "columns": [],
        },
        {
            "index": 2,
            "layout": "bullets",
            "title": "为什么这节课重要",
            "subtitle": "",
            "bullets": ["岗位面试常考", "承接后续 numpy 广播"],
            "notes": "留 30 秒抢答",
            "key_points": ["面试高频"],
            "citations": [{"title": "案例A", "url": "https://example.com/a"}],
            "columns": [],
        },
        {
            "index": 3,
            "layout": "two_column",
            "title": "切片 vs 元组",
            "subtitle": "",
            "bullets": [],
            "notes": "",
            "key_points": [],
            "citations": [],
            "columns": [
                {"title": "切片", "bullets": ["支持负索引", "步长"]},
                {"title": "元组", "bullets": ["不可变", "可作为 key"]},
            ],
        },
        {
            "index": 4,
            "layout": "callout",
            "title": "重点提醒",
            "subtitle": "",
            "bullets": ["切片返回浅拷贝"],
            "notes": "",
            "key_points": [],
            "citations": [],
            "columns": [],
        },
        {
            "index": 5,
            "layout": "summary",
            "title": "小结",
            "subtitle": "",
            "bullets": ["切片语法", "元组特性"],
            "notes": "",
            "key_points": [],
            "citations": [],
            "columns": [],
        },
    ],
    "sources": [{"title": "案例A", "url": "https://example.com/a", "snippet": "..."}],
}


@pytest.mark.parametrize("template_id", ["ai_tech", "business_plan"])
def test_slide_deck_pptx_is_a_valid_template_deck(template_id: str) -> None:
    exported = SlideDeckPptxSkill().run({**SAMPLE_DECK, "template_id": template_id})

    assert exported.extension == "pptx"
    assert exported.media_type.endswith("presentationml.presentation")
    # PPTX files are ZIP archives; the first two bytes must be the ZIP magic.
    assert exported.content[:2] == b"PK"

    with zipfile.ZipFile(io.BytesIO(exported.content)) as archive:
        slide_files = [
            name
            for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
        assert all(not _empty_placeholder_ids(archive.read(name)) for name in slide_files)
    # The declared title slide is the cover; the exporter must not duplicate it.
    assert len(slide_files) == len(SAMPLE_DECK["slides"])

    presentation = Presentation(io.BytesIO(exported.content))
    assert len(presentation.slides) == len(SAMPLE_DECK["slides"])
    assert round(presentation.slide_width / 914400, 2) == 13.33
    assert round(presentation.slide_height / 914400, 2) == 7.5
    text = _presentation_text(presentation)
    assert "Python 切片与元组" in text
    assert "第一PPT" not in text
    assert "物联网产品" not in text
    assert "AI basic knowledge training" not in text
    assert "www.1ppt.com" not in text


def test_slide_deck_without_title_gets_one_synthesised_cover() -> None:
    data = {**SAMPLE_DECK, "slides": SAMPLE_DECK["slides"][1:]}

    exported = SlideDeckPptxSkill().run(data)
    presentation = Presentation(io.BytesIO(exported.content))

    assert len(presentation.slides) == len(data["slides"]) + 1
    assert "Python 切片与元组" in _presentation_text(presentation)


def test_unknown_template_id_falls_back_to_ai_tech() -> None:
    exported = SlideDeckPptxSkill().run({**SAMPLE_DECK, "template_id": "../../other"})
    presentation = Presentation(io.BytesIO(exported.content))

    assert presentation.slide_layouts[0].name == "Title Slide"


def test_slide_deck_json_preserves_optional_template_id() -> None:
    result = SlideDeckJsonSkill().run({**SAMPLE_DECK, "template_id": "business_plan"})

    assert result["template_id"] == "business_plan"


def test_template_catalog_and_manifests_are_valid() -> None:
    specs = template_catalog()

    assert [spec.id for spec in specs] == ["ai_tech", "business_plan"]
    assert all(spec.license_scope == "development_only" for spec in specs)
    assert find_explicit_template_id("请使用商业计划书模板生成课件") == "business_plan"
    assert find_explicit_template_id("使用科技蓝模板") == "ai_tech"
    for spec in specs:
        validate_template_manifest(spec)


def test_frame_planner_avoids_adjacent_repetition_when_an_alternative_fits() -> None:
    data = {
        **SAMPLE_DECK,
        "slides": [
            SAMPLE_DECK["slides"][0],
            *[
                {
                    **SAMPLE_DECK["slides"][1],
                    "index": index,
                    "title": f"要点 {index}",
                }
                for index in range(2, 7)
            ],
        ],
    }

    planned = plan_template_slides(get_template("ai_tech"), data)
    content_frame_ids = [item.frame.id for item in planned[1:]]

    assert all(
        left != right for left, right in zip(content_frame_ids, content_frame_ids[1:])
    )


def _presentation_text(presentation) -> str:
    def iter_shapes(shapes):
        for shape in shapes:
            yield shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from iter_shapes(shape.shapes)

    return "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in iter_shapes(slide.shapes)
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    )


def _empty_placeholder_ids(slide_xml: bytes) -> list[str]:
    root = etree.fromstring(slide_xml)
    namespaces = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }
    placeholders = root.xpath(
        ".//p:sp[p:nvSpPr/p:nvPr/p:ph and not(.//a:t[normalize-space()])]",
        namespaces=namespaces,
    )
    return [
        str(shape.xpath("string(p:nvSpPr/p:cNvPr/@id)", namespaces=namespaces))
        for shape in placeholders
    ]
